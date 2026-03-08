#!/usr/bin/env python3
"""
Build a PPTX deck from full-bleed slide images.

Use when images already contain background + text.
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation


def iter_slide_images(indir: Path) -> list[Path]:
    imgs = []
    for p in indir.glob("slide-*.png"):
        m = re.match(r"slide-(\d+)\.png$", p.name)
        if not m:
            continue
        imgs.append((int(m.group(1)), p))
    imgs.sort(key=lambda x: x[0])
    return [p for _, p in imgs]


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--indir", required=True, help="Directory containing slide-XX.png files")
    ap.add_argument("--out", required=True, help="Output pptx path")
    args = ap.parse_args()

    indir = Path(args.indir).resolve()
    out = Path(args.out).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    images = iter_slide_images(indir)
    if not images:
        raise SystemExit(f"No slide-XX.png files found in {indir}")

    prs = Presentation()
    # 16:9 widescreen (13.333 x 7.5 in)
    prs.slide_width = int(13.333 * 914400)
    prs.slide_height = int(7.5 * 914400)

    blank_layout = prs.slide_layouts[6]
    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(str(out))
    print(out)
    print(f"slides={len(images)}")


if __name__ == "__main__":
    main()

